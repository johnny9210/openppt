"""
PPTX Exporter
Converts pipeline output to .pptx file.

Two modes:
  1. Image-based (export_pptx_from_images): Uses captured Preview screenshots as slides
  2. Text-based (export_pptx): Fallback using structured text content with styling
"""

import base64
import io
import logging
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

logger = logging.getLogger(__name__)

# Slide dimensions: 16:9 widescreen
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Design system defaults (matches Preview THEME) ──────
DEFAULT_THEME = {
    "primary_color": "#6366F1",
    "accent_color": "#818CF8",
    "background": "#F5F7FA",
    "text_color": "#1A202C",
}

# Extended colors (same as code_assembly.py THEME)
TEXT_SECONDARY = "#64748B"
CARD_BG = "#FFFFFF"
CARD_BORDER = "#E2E8F0"
DIVIDER = "#E2E8F0"
RED = "#E53E3E"
YELLOW = "#F59E0B"
GREEN = "#38A169"
SHADOW_COLOR = "#C0C4CC"
SHADOW_OFFSET = 0.06  # inches offset for card shadow


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert #RRGGBB to RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_slide_background(slide, hex_color: str) -> None:
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex_color)


def _add_shape(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str = CARD_BG,
    border_color: str | None = CARD_BORDER,
    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
) -> Any:
    """Add a shape (card, circle, bar) to a slide."""
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = _hex_to_rgb(border_color)
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def _add_accent_bar(slide, left: float, top: float, color: str, width: float = 0.5, height: float = 0.04) -> None:
    """Add a small colored accent bar (like the underline under titles in Preview)."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(color)
    bar.line.fill.background()


def _add_circle(
    slide,
    left: float,
    top: float,
    size: float,
    fill_color: str,
    text: str = "",
    font_size: int = 14,
) -> None:
    """Add a colored circle (icon badge) with optional text."""
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(size), Inches(size),
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = _hex_to_rgb(fill_color)
    circle.line.fill.background()
    if text:
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int = 18,
    color: str = "#1A202C",
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str | None = None,
) -> Any:
    """Add a text box to a slide and return the text frame."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.color.rgb = _hex_to_rgb(color)
    p.font.bold = bold
    p.alignment = alignment
    if font_name:
        p.font.name = font_name
    return tf


def _add_paragraph(
    text_frame,
    text: str,
    font_size: int = 14,
    color: str = TEXT_SECONDARY,
    bold: bool = False,
    space_before: int = 6,
) -> None:
    """Append a paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.color.rgb = _hex_to_rgb(color)
    p.font.bold = bold
    p.space_before = Pt(space_before)


def _add_card_with_shadow(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str = CARD_BG,
    border_color: str | None = CARD_BORDER,
    shadow_color: str = SHADOW_COLOR,
) -> Any:
    """Add a card with a subtle shadow effect (offset darker rectangle behind the white card)."""
    # Shadow rectangle (slightly offset down-right)
    _add_shape(
        slide, left + SHADOW_OFFSET, top + SHADOW_OFFSET, width, height,
        fill_color=shadow_color, border_color=None,
        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    # Main card on top
    return _add_shape(slide, left, top, width, height, fill_color=fill_color, border_color=border_color)


def _add_slide_number(slide, number: int, total: int) -> None:
    """Add a small slide number indicator (e.g. '3 / 10') at bottom-right of the slide."""
    _add_textbox(
        slide, 11.0, 7.05, 1.8, 0.3,
        f"{number} / {total}",
        font_size=9, color=TEXT_SECONDARY, bold=False,
        alignment=PP_ALIGN.RIGHT,
    )


# ── Slide type renderers (matching Preview design system) ──────

def _render_cover(slide, content: dict, theme: dict) -> None:
    """Render cover slide — left-aligned title with subtitle and presenter."""
    primary = theme.get("primary_color", "#6366F1")
    accent = theme.get("accent_color", "#818CF8")
    text_color = theme.get("text_color", "#1A202C")

    # Decorative colored strip behind the title area (gradient-like effect)
    _add_shape(
        slide, 0, 1.6, 0.15, 4.2,
        fill_color=primary, border_color=None,
        shape_type=MSO_SHAPE.RECTANGLE,
    )

    # Right-side decorative accent bar
    _add_shape(
        slide, 12.5, 0.8, 0.08, 5.8,
        fill_color=accent, border_color=None,
        shape_type=MSO_SHAPE.RECTANGLE,
    )

    # Title (left-aligned, large, better vertical centering)
    title = content.get("title", "")
    title_size = 32 if len(title) > 30 else 44
    _add_textbox(
        slide, 1.0, 2.2, 8.0, 2.0,
        title,
        font_size=title_size, color=text_color, bold=True,
    )

    # Accent bar under title
    _add_accent_bar(slide, 1.0, 4.0, primary, width=0.6, height=0.05)

    # Subtitle
    if content.get("subtitle"):
        _add_textbox(
            slide, 1.0, 4.3, 8.0, 0.8,
            content["subtitle"],
            font_size=18, color=TEXT_SECONDARY,
        )

    # Presenter + Date (bottom left)
    parts = []
    if content.get("presenter"):
        parts.append(content["presenter"])
    if content.get("date"):
        parts.append(content["date"])
    if parts:
        _add_textbox(
            slide, 1.0, 6.2, 8.0, 0.4,
            "  |  ".join(parts),
            font_size=13, color=TEXT_SECONDARY,
        )


def _render_table_of_contents(slide, content: dict, theme: dict) -> None:
    """Render table of contents with numbered items and accent colors."""
    primary = theme.get("primary_color", "#6366F1")
    accent = theme.get("accent_color", "#818CF8")
    text_color = theme.get("text_color", "#1A202C")

    # Title
    _add_textbox(
        slide, 1.0, 0.5, 11.3, 0.7,
        content.get("title", "Contents"),
        font_size=32, color=text_color, bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    _add_accent_bar(slide, 6.4, 1.15, primary, width=0.5, height=0.04)

    items = content.get("items", [])
    y = 1.6
    for i, item in enumerate(items):
        number = item.get("number", str(i + 1))
        title = item.get("title", "")
        desc = item.get("description", "")

        # Number circle
        circle_color = primary if i % 2 == 0 else accent
        _add_circle(slide, 1.5, y, 0.45, circle_color, str(number), font_size=14)

        # Title
        _add_textbox(
            slide, 2.2, y, 9.0, 0.35,
            title,
            font_size=20, color=text_color, bold=True,
        )
        # Description
        if desc:
            _add_textbox(
                slide, 2.2, y + 0.35, 9.0, 0.3,
                desc,
                font_size=14, color=TEXT_SECONDARY,
            )
        y += 0.85


def _generate_chart_palette(primary_hex: str, accent_hex: str, count: int) -> list[RGBColor]:
    """Generate *count* RGBColor values for chart slices / points.

    Interpolates between primary and accent so pie charts with many slices
    still look cohesive with the slide theme.
    """
    def _lerp_hex(hex_a: str, hex_b: str, t: float) -> RGBColor:
        a = hex_a.lstrip("#")
        b = hex_b.lstrip("#")
        r = int(int(a[0:2], 16) * (1 - t) + int(b[0:2], 16) * t)
        g = int(int(a[2:4], 16) * (1 - t) + int(b[2:4], 16) * t)
        b_val = int(int(a[4:6], 16) * (1 - t) + int(b[4:6], 16) * t)
        return RGBColor(min(r, 255), min(g, 255), min(b_val, 255))

    if count <= 0:
        return []
    base = [_hex_to_rgb(primary_hex), _hex_to_rgb(accent_hex)]
    if count <= 2:
        return base[:count]
    palette: list[RGBColor] = []
    for i in range(count):
        t = i / max(count - 1, 1)
        palette.append(_lerp_hex(primary_hex, accent_hex, t))
    return palette


def _render_data_visualization(slide, content: dict, theme: dict) -> None:
    """Render data visualization slide with a native PowerPoint chart.

    Supports chart_type: bar, line, pie, area.  Unsupported types (radar,
    scatter, funnel, etc.) gracefully fall back to a bar chart.  When the
    data list is empty or contains no usable numeric values the function
    falls back to a text-only placeholder so the slide is never blank.
    """
    primary = theme.get("primary_color", "#6366F1")
    accent = theme.get("accent_color", "#818CF8")
    text_color = theme.get("text_color", "#1A202C")

    # ── Title ───────────────────────────────────────────────
    _add_textbox(
        slide, 0.8, 0.4, 11.7, 0.7,
        content.get("title", ""),
        font_size=28, color=text_color, bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    _add_accent_bar(slide, 6.4, 1.05, primary, width=0.5, height=0.04)

    if content.get("description"):
        _add_textbox(
            slide, 0.8, 1.3, 11.7, 0.5,
            content["description"],
            font_size=14, color=TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER,
        )

    # ── Prepare data ────────────────────────────────────────
    raw_data = content.get("data") or []
    categories: list[str] = []
    values: list[float] = []
    for item in raw_data:
        name = str(item.get("name", ""))
        raw_value = item.get("value")
        try:
            numeric = float(raw_value)
        except (TypeError, ValueError):
            continue  # skip non-numeric entries
        categories.append(name)
        values.append(numeric)

    # ── Chart type mapping ──────────────────────────────────
    CHART_TYPE_MAP = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "area": XL_CHART_TYPE.AREA,
    }
    requested_type = (content.get("chart_type") or "bar").lower().strip()
    xl_chart_type = CHART_TYPE_MAP.get(requested_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    # ── Build chart or fallback ─────────────────────────────
    if categories and values:
        chart_data = CategoryChartData()
        chart_data.categories = categories
        series_label = content.get("title", "Value") or "Value"
        chart_data.add_series(series_label, values)

        # Chart placement — centred card-like area
        chart_left = Inches(1.8)
        chart_top = Inches(2.0)
        chart_width = Inches(9.7)
        chart_height = Inches(3.8)

        graphic_frame = slide.shapes.add_chart(
            xl_chart_type, chart_left, chart_top, chart_width, chart_height,
            chart_data,
        )
        chart = graphic_frame.chart

        # ── Style the chart ─────────────────────────────────
        chart.has_title = False  # title is already rendered as a textbox

        # Legend — subtle, bottom placement for bar/line/area; right for pie
        if xl_chart_type == XL_CHART_TYPE.PIE:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(9)
            chart.legend.font.color.rgb = _hex_to_rgb(TEXT_SECONDARY)
        elif len(categories) > 1:
            chart.has_legend = False  # single-series; legend is noise
        else:
            chart.has_legend = False

        # Axes — clean styling (pie charts have no axes)
        if xl_chart_type != XL_CHART_TYPE.PIE:
            # Category axis
            cat_axis = chart.category_axis
            cat_axis.has_major_gridlines = False
            cat_axis.has_minor_gridlines = False
            cat_axis.tick_labels.font.size = Pt(9)
            cat_axis.tick_labels.font.color.rgb = _hex_to_rgb(TEXT_SECONDARY)
            cat_axis.format.line.fill.background()  # hide axis line

            # Value axis
            val_axis = chart.value_axis
            val_axis.has_major_gridlines = True
            val_axis.major_gridlines.format.line.color.rgb = _hex_to_rgb(DIVIDER)
            val_axis.major_gridlines.format.line.width = Pt(0.4)
            val_axis.has_minor_gridlines = False
            val_axis.tick_labels.font.size = Pt(9)
            val_axis.tick_labels.font.color.rgb = _hex_to_rgb(TEXT_SECONDARY)
            val_axis.format.line.fill.background()  # hide axis line

        # Series colouring
        primary_rgb = _hex_to_rgb(primary)

        if xl_chart_type == XL_CHART_TYPE.PIE:
            # Colour each slice with a cohesive gradient palette
            plot = chart.plots[0]
            series = plot.series[0]
            _extended = _generate_chart_palette(primary, accent, len(categories))
            for idx in range(len(categories)):
                point = series.points[idx]
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = _extended[idx]
        else:
            # Bar / Line / Area — colour the single series
            series = chart.series[0]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = primary_rgb
            if xl_chart_type == XL_CHART_TYPE.LINE_MARKERS:
                series.format.line.color.rgb = primary_rgb
                series.format.line.width = Pt(2.5)
                series.smooth = True
    else:
        # Fallback: no usable data — render a placeholder card
        _add_shape(slide, 2.0, 2.2, 9.3, 3.5)
        _add_textbox(
            slide, 2.5, 3.5, 8.3, 0.5,
            "No chart data available",
            font_size=16, color=TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER,
        )

    # ── Insight ─────────────────────────────────────────────
    if content.get("insight"):
        _add_textbox(
            slide, 0.8, 6.2, 11.7, 0.6,
            f"💡 {content['insight']}",
            font_size=14, color=primary, bold=True,
            alignment=PP_ALIGN.CENTER,
        )


def _render_key_points(slide, content: dict, theme: dict) -> None:
    """Render key points with cards and icon badges (matching Preview)."""
    primary = theme.get("primary_color", "#6366F1")
    accent = theme.get("accent_color", "#818CF8")
    text_color = theme.get("text_color", "#1A202C")

    # Title
    _add_textbox(
        slide, 1.0, 0.4, 11.3, 0.7,
        content.get("title", ""),
        font_size=32, color=text_color, bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    _add_accent_bar(slide, 6.4, 1.05, primary, width=0.5, height=0.04)

    if content.get("description"):
        _add_textbox(
            slide, 1.0, 1.25, 11.3, 0.4,
            content["description"],
            font_size=14, color=TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER,
        )

    points = content.get("points", content.get("items", []))
    if not points:
        return

    # Layout: match Preview — centered icon on top, text below
    cols = min(len(points), 3)
    gap = 0.35
    card_w = (11.0 - (cols - 1) * gap) / cols
    start_x = (13.333 - (card_w * cols + gap * (cols - 1))) / 2
    start_y = 1.9
    card_h = 3.2

    for i, pt in enumerate(points):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + 0.3)

        # Card background with shadow
        _add_card_with_shadow(slide, x, y, card_w, card_h)

        # Icon circle (centered at top of card)
        emoji = pt.get("emoji", "●")
        circle_color = primary if i % 2 == 0 else accent
        icon_size = 0.65
        cx = x + (card_w - icon_size) / 2
        _add_circle(slide, cx, y + 0.3, icon_size, circle_color, emoji, font_size=16)

        # Title (centered below icon)
        _add_textbox(
            slide, x + 0.2, y + 1.15, card_w - 0.4, 0.4,
            pt.get("title", ""),
            font_size=17, color=text_color, bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # Description (centered below title)
        if pt.get("description"):
            _add_textbox(
                slide, x + 0.2, y + 1.6, card_w - 0.4, 1.0,
                pt["description"],
                font_size=12, color=TEXT_SECONDARY,
                alignment=PP_ALIGN.CENTER,
            )

        # Metric (centered at bottom)
        if pt.get("metric"):
            _add_textbox(
                slide, x + 0.2, y + 2.65, card_w - 0.4, 0.35,
                pt["metric"],
                font_size=14, color=primary, bold=True,
                alignment=PP_ALIGN.CENTER,
            )


def _render_risk_analysis(slide, content: dict, theme: dict) -> None:
    """Render risk analysis with color-coded cards."""
    primary = theme.get("primary_color", "#6366F1")
    text_color = theme.get("text_color", "#1A202C")

    LEVEL_COLORS = {"high": RED, "medium": YELLOW, "low": GREEN}

    # Title
    _add_textbox(
        slide, 1.0, 0.4, 11.3, 0.7,
        content.get("title", ""),
        font_size=28, color=text_color, bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    _add_accent_bar(slide, 6.4, 1.05, primary, width=0.5, height=0.04)

    if content.get("description"):
        _add_textbox(
            slide, 1.0, 1.25, 11.3, 0.4,
            content["description"],
            font_size=14, color=TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER,
        )

    risks = content.get("risks", [])
    y = 1.9
    for risk in risks:
        level = risk.get("level", "medium")
        level_color = LEVEL_COLORS.get(level, primary)
        title = risk.get("title", "")
        desc = risk.get("description", "")
        mitigation = risk.get("mitigation", "")

        card_h = 0.8 + (0.3 if desc else 0) + (0.3 if mitigation else 0)
        _add_card_with_shadow(slide, 1.5, y, 10.3, card_h)

        # Level badge (colored bar on left edge)
        _add_shape(
            slide, 1.5, y, 0.06, card_h,
            fill_color=level_color, border_color=None,
            shape_type=MSO_SHAPE.RECTANGLE,
        )

        # Level + Title
        _add_textbox(
            slide, 1.8, y + 0.1, 2.0, 0.35,
            f"[{level.upper()}]",
            font_size=12, color=level_color, bold=True,
        )
        _add_textbox(
            slide, 3.5, y + 0.1, 8.0, 0.35,
            title,
            font_size=17, color=text_color, bold=True,
        )

        inner_y = y + 0.45
        if desc:
            _add_textbox(
                slide, 1.8, inner_y, 9.8, 0.3,
                desc,
                font_size=13, color=TEXT_SECONDARY,
            )
            inner_y += 0.3
        if mitigation:
            _add_textbox(
                slide, 1.8, inner_y, 9.8, 0.3,
                f"→ {mitigation}",
                font_size=12, color=primary,
            )

        y += card_h + 0.15


def _render_action_plan(slide, content: dict, theme: dict) -> None:
    """Render action plan with phase cards."""
    primary = theme.get("primary_color", "#6366F1")
    accent = theme.get("accent_color", "#818CF8")
    text_color = theme.get("text_color", "#1A202C")

    # Title
    _add_textbox(
        slide, 1.0, 0.4, 11.3, 0.7,
        content.get("title", ""),
        font_size=28, color=text_color, bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    _add_accent_bar(slide, 6.4, 1.05, primary, width=0.5, height=0.04)

    if content.get("description"):
        _add_textbox(
            slide, 1.0, 1.25, 11.3, 0.4,
            content["description"],
            font_size=14, color=TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER,
        )

    actions = content.get("actions", [])
    y = 1.9
    for i, action in enumerate(actions):
        phase = action.get("phase", "")
        title = action.get("title", "")
        period = action.get("period", "")
        tasks = action.get("tasks", [])

        card_h = 0.7 + len(tasks) * 0.3
        _add_card_with_shadow(slide, 1.5, y, 10.3, card_h)

        # Phase circle
        circle_color = primary if i % 2 == 0 else accent
        _add_circle(slide, 1.7, y + 0.15, 0.45, circle_color, phase or str(i + 1), font_size=11)

        # Title + period
        header = title
        if period:
            header += f"  ({period})"
        _add_textbox(
            slide, 2.4, y + 0.15, 9.0, 0.35,
            header,
            font_size=17, color=text_color, bold=True,
        )

        task_y = y + 0.5
        for task in tasks:
            _add_textbox(
                slide, 2.6, task_y, 8.8, 0.25,
                f"•  {task}",
                font_size=13, color=TEXT_SECONDARY,
            )
            task_y += 0.28

        y += card_h + 0.15


def _render_hero(slide, content: dict, theme: dict) -> None:
    """Render hero slide — large centered title."""
    primary = theme.get("primary_color", "#6366F1")
    text_color = theme.get("text_color", "#1A202C")

    _add_textbox(
        slide, 1.5, 2.0, 10.3, 2.0,
        content.get("title", ""),
        font_size=48, color=text_color, bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    if content.get("subtitle"):
        _add_textbox(
            slide, 2.0, 4.2, 9.3, 1.0,
            content["subtitle"],
            font_size=20, color=TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER,
        )


def _render_quote(slide, content: dict, theme: dict) -> None:
    """Render quote slide with decorative accent."""
    primary = theme.get("primary_color", "#6366F1")
    text_color = theme.get("text_color", "#1A202C")

    # Accent bar
    _add_accent_bar(slide, 6.6, 1.5, primary, width=0.05, height=0.6)

    # Quote text
    _add_textbox(
        slide, 2.0, 2.5, 9.3, 2.0,
        content.get("quote", ""),
        font_size=28, color=text_color, bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    if content.get("attribution"):
        _add_textbox(
            slide, 2.0, 4.8, 9.3, 0.4,
            f"— {content['attribution']}",
            font_size=16, color=TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER,
        )

    if content.get("context"):
        _add_textbox(
            slide, 2.0, 5.3, 9.3, 0.4,
            content["context"],
            font_size=14, color=TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER,
        )


def _render_icon_grid(slide, content: dict, theme: dict) -> None:
    """Render icon grid with cards."""
    primary = theme.get("primary_color", "#6366F1")
    accent = theme.get("accent_color", "#818CF8")
    text_color = theme.get("text_color", "#1A202C")

    # Title
    _add_textbox(
        slide, 1.0, 0.4, 11.3, 0.7,
        content.get("title", ""),
        font_size=32, color=text_color, bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    _add_accent_bar(slide, 6.4, 1.05, primary, width=0.5, height=0.04)

    if content.get("description"):
        _add_textbox(
            slide, 1.0, 1.25, 11.3, 0.4,
            content["description"],
            font_size=14, color=TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER,
        )

    items = content.get("items", [])
    if not items:
        return

    cols = 2 if len(items) <= 4 else 3
    card_w = (10.5 - (cols - 1) * 0.25) / cols
    start_x = (13.333 - (card_w * cols + 0.25 * (cols - 1))) / 2
    start_y = 1.9

    for i, item in enumerate(items):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + 0.25)
        y = start_y + row * 1.8

        _add_card_with_shadow(slide, x, y, card_w, 1.6)

        # Icon circle (centered)
        emoji = item.get("emoji", "●")
        circle_color = primary if i % 2 == 0 else accent
        cx = x + (card_w - 0.55) / 2
        _add_circle(slide, cx, y + 0.15, 0.55, circle_color, emoji, font_size=14)

        # Label
        _add_textbox(
            slide, x + 0.2, y + 0.8, card_w - 0.4, 0.3,
            item.get("label", item.get("title", "")),
            font_size=16, color=text_color, bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # Description
        if item.get("description"):
            _add_textbox(
                slide, x + 0.2, y + 1.1, card_w - 0.4, 0.4,
                item["description"],
                font_size=12, color=TEXT_SECONDARY,
                alignment=PP_ALIGN.CENTER,
            )


def _render_three_column(slide, content: dict, theme: dict) -> None:
    """Render three-column layout with cards."""
    primary = theme.get("primary_color", "#6366F1")
    accent = theme.get("accent_color", "#818CF8")
    text_color = theme.get("text_color", "#1A202C")

    _add_textbox(
        slide, 1.0, 0.4, 11.3, 0.7,
        content.get("title", ""),
        font_size=32, color=text_color, bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    _add_accent_bar(slide, 6.4, 1.05, primary, width=0.5, height=0.04)

    if content.get("description"):
        _add_textbox(
            slide, 1.0, 1.25, 11.3, 0.4,
            content["description"],
            font_size=14, color=TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER,
        )

    columns = content.get("columns", [])
    if not columns:
        return

    cols = len(columns)
    card_w = (11.0 - (cols - 1) * 0.3) / cols
    start_x = (13.333 - (card_w * cols + 0.3 * (cols - 1))) / 2
    y = 1.9

    colors = [primary, accent, GREEN]
    for i, col in enumerate(columns):
        x = start_x + i * (card_w + 0.3)
        _add_card_with_shadow(slide, x, y, card_w, 4.5)

        # Icon circle
        emoji = col.get("emoji", "●")
        c = colors[i % len(colors)]
        cx = x + (card_w - 0.55) / 2
        _add_circle(slide, cx, y + 0.25, 0.55, c, emoji, font_size=14)

        # Title
        _add_textbox(
            slide, x + 0.2, y + 0.95, card_w - 0.4, 0.35,
            col.get("title", ""),
            font_size=17, color=text_color, bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # Description
        if col.get("description"):
            _add_textbox(
                slide, x + 0.2, y + 1.4, card_w - 0.4, 1.5,
                col["description"],
                font_size=13, color=TEXT_SECONDARY,
                alignment=PP_ALIGN.CENTER,
            )

        # Metric
        if col.get("metric"):
            _add_textbox(
                slide, x + 0.2, y + 3.5, card_w - 0.4, 0.5,
                col["metric"],
                font_size=20, color=primary, bold=True,
                alignment=PP_ALIGN.CENTER,
            )


def _render_comparison(slide, content: dict, theme: dict) -> None:
    """Render comparison slide with two side-by-side cards."""
    primary = theme.get("primary_color", "#6366F1")
    text_color = theme.get("text_color", "#1A202C")

    _add_textbox(
        slide, 1.0, 0.4, 11.3, 0.7,
        content.get("title", ""),
        font_size=32, color=text_color, bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    left = content.get("left", {"label": "", "items": []})
    right = content.get("right", {"label": "", "items": []})

    card_w = 5.3
    card_h = 5.0

    for side_idx, (data, color, icon) in enumerate([
        (left, RED, "✕"),
        (right, GREEN, "✓"),
    ]):
        x = 1.2 + side_idx * (card_w + 0.3)
        y = 1.5
        _add_card_with_shadow(slide, x, y, card_w, card_h)

        # Header circle + label
        _add_circle(slide, x + 0.3, y + 0.3, 0.4, color, icon, font_size=14)
        _add_textbox(
            slide, x + 0.85, y + 0.3, card_w - 1.2, 0.4,
            data.get("label", ""),
            font_size=18, color=text_color, bold=True,
        )

        # Items
        items = data.get("items", [])
        item_y = y + 0.9
        for item in items:
            mark_color = RED if color == RED else GREEN
            _add_textbox(
                slide, x + 0.4, item_y, 0.3, 0.3,
                icon,
                font_size=13, color=mark_color,
            )
            _add_textbox(
                slide, x + 0.7, item_y, card_w - 1.0, 0.3,
                str(item),
                font_size=14, color=TEXT_SECONDARY,
            )
            item_y += 0.38


def _render_process_flow(slide, content: dict, theme: dict) -> None:
    """Render process flow with step cards and arrows."""
    primary = theme.get("primary_color", "#6366F1")
    accent = theme.get("accent_color", "#818CF8")
    text_color = theme.get("text_color", "#1A202C")

    _add_textbox(
        slide, 1.0, 0.4, 11.3, 0.7,
        content.get("title", ""),
        font_size=32, color=text_color, bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    _add_accent_bar(slide, 6.4, 1.05, primary, width=0.5, height=0.04)

    if content.get("description"):
        _add_textbox(
            slide, 1.0, 1.25, 11.3, 0.4,
            content["description"],
            font_size=14, color=TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER,
        )

    steps = content.get("steps", [])
    if not steps:
        return

    n = len(steps)
    arrow_w = 0.25
    total_arrows = max(0, n - 1) * arrow_w
    card_w = min(1.8, (11.0 - total_arrows) / n)
    total_w = n * card_w + total_arrows
    start_x = (13.333 - total_w) / 2
    y = 2.0

    for i, step in enumerate(steps):
        x = start_x + i * (card_w + arrow_w)
        _add_card_with_shadow(slide, x, y, card_w, 3.5)

        # Step number/emoji circle
        emoji = step.get("emoji", str(i + 1))
        c = primary if i % 2 == 0 else accent
        cx = x + (card_w - 0.45) / 2
        _add_circle(slide, cx, y + 0.2, 0.45, c, emoji, font_size=12)

        # Title
        _add_textbox(
            slide, x + 0.1, y + 0.8, card_w - 0.2, 0.5,
            step.get("title", ""),
            font_size=14, color=text_color, bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # Description
        if step.get("description"):
            _add_textbox(
                slide, x + 0.1, y + 1.3, card_w - 0.2, 1.5,
                step["description"],
                font_size=11, color=TEXT_SECONDARY,
                alignment=PP_ALIGN.CENTER,
            )

        # Arrow between steps
        if i < n - 1:
            arrow_x = x + card_w + 0.02
            _add_textbox(
                slide, arrow_x, y + 1.3, arrow_w, 0.4,
                "→",
                font_size=20, color=primary, bold=True,
                alignment=PP_ALIGN.CENTER,
            )


def _render_timeline(slide, content: dict, theme: dict) -> None:
    """Render timeline with events along a horizontal line."""
    primary = theme.get("primary_color", "#6366F1")
    accent = theme.get("accent_color", "#818CF8")
    text_color = theme.get("text_color", "#1A202C")

    _add_textbox(
        slide, 1.0, 0.4, 11.3, 0.7,
        content.get("title", ""),
        font_size=32, color=text_color, bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    _add_accent_bar(slide, 6.4, 1.05, primary, width=0.5, height=0.04)

    if content.get("description"):
        _add_textbox(
            slide, 1.0, 1.25, 11.3, 0.4,
            content["description"],
            font_size=14, color=TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER,
        )

    events = content.get("events", [])
    if not events:
        return

    # Timeline line
    _add_shape(
        slide, 1.5, 3.4, 10.3, 0.03,
        fill_color=primary, border_color=None,
        shape_type=MSO_SHAPE.RECTANGLE,
    )

    n = len(events)
    spacing = 10.3 / max(n, 1)
    start_x = 1.5

    for i, evt in enumerate(events):
        x = start_x + i * spacing + spacing / 2 - 0.6
        c = primary if i % 2 == 0 else accent

        # Time label
        _add_textbox(
            slide, x, 2.3, 1.2, 0.3,
            evt.get("time", ""),
            font_size=11, color=primary, bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # Circle on timeline
        _add_circle(slide, x + 0.35, 3.15, 0.4, c, evt.get("emoji", "●"), font_size=12)

        # Card below
        _add_card_with_shadow(slide, x - 0.1, 3.7, 1.4, 1.5)
        _add_textbox(
            slide, x, 3.8, 1.2, 0.3,
            evt.get("title", ""),
            font_size=12, color=text_color, bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        if evt.get("description"):
            _add_textbox(
                slide, x, 4.15, 1.2, 0.8,
                evt["description"],
                font_size=10, color=TEXT_SECONDARY,
                alignment=PP_ALIGN.CENTER,
            )


def _render_summary(slide, content: dict, theme: dict) -> None:
    """Render summary slide with numbered points in cards."""
    primary = theme.get("primary_color", "#6366F1")
    text_color = theme.get("text_color", "#1A202C")

    _add_textbox(
        slide, 1.0, 0.4, 11.3, 0.7,
        content.get("title", ""),
        font_size=32, color=text_color, bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    _add_accent_bar(slide, 6.4, 1.05, primary, width=0.5, height=0.04)

    points = content.get("points", [])
    y = 1.5
    for i, pt in enumerate(points):
        _add_card_with_shadow(slide, 2.5, y, 8.3, 0.8)

        # Number circle
        _add_circle(slide, 2.7, y + 0.15, 0.5, primary, str(pt.get("number", i + 1)), font_size=14)

        _add_textbox(
            slide, 3.4, y + 0.1, 7.0, 0.35,
            pt.get("title", ""),
            font_size=17, color=text_color, bold=True,
        )
        if pt.get("description"):
            _add_textbox(
                slide, 3.4, y + 0.45, 7.0, 0.3,
                pt["description"],
                font_size=13, color=TEXT_SECONDARY,
            )
        y += 0.95


def _render_closing(slide, content: dict, theme: dict) -> None:
    """Render closing slide with resources."""
    primary = theme.get("primary_color", "#6366F1")
    accent = theme.get("accent_color", "#818CF8")
    text_color = theme.get("text_color", "#1A202C")

    # Decorative accent bar at top center
    _add_accent_bar(slide, 6.15, 1.1, primary, width=1.0, height=0.05)

    # Title or "Thank You" fallback
    title = content.get("title", "")
    if not title:
        title = "Thank You"
    _add_textbox(
        slide, 1.5, 1.5, 10.3, 1.0,
        title,
        font_size=36, color=text_color, bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    if content.get("message"):
        _add_textbox(
            slide, 2.5, 2.8, 8.3, 1.0,
            content["message"],
            font_size=16, color=TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER,
        )

    resources = content.get("resources", [])
    if resources:
        n = len(resources)
        card_w = min(2.0, (10.0 - (n - 1) * 0.25) / n)
        total_w = n * card_w + (n - 1) * 0.25
        start_x = (13.333 - total_w) / 2
        y = 4.2

        for i, res in enumerate(resources):
            x = start_x + i * (card_w + 0.25)
            _add_card_with_shadow(slide, x, y, card_w, 1.5)

            # Emoji
            _add_textbox(
                slide, x, y + 0.15, card_w, 0.45,
                res.get("emoji", ""),
                font_size=24, color=text_color,
                alignment=PP_ALIGN.CENTER,
            )
            _add_textbox(
                slide, x + 0.1, y + 0.65, card_w - 0.2, 0.35,
                res.get("label", ""),
                font_size=13, color=text_color, bold=True,
                alignment=PP_ALIGN.CENTER,
            )
            if res.get("description"):
                _add_textbox(
                    slide, x + 0.1, y + 1.0, card_w - 0.2, 0.4,
                    res["description"],
                    font_size=10, color=TEXT_SECONDARY,
                    alignment=PP_ALIGN.CENTER,
                )

    # Bottom decorative accent bar
    _add_accent_bar(slide, 5.65, 6.5, accent, width=2.0, height=0.04)


# Renderer dispatch
RENDERERS = {
    "cover": _render_cover,
    "table_of_contents": _render_table_of_contents,
    "data_visualization": _render_data_visualization,
    "key_points": _render_key_points,
    "risk_analysis": _render_risk_analysis,
    "action_plan": _render_action_plan,
    "hero": _render_hero,
    "quote": _render_quote,
    "icon_grid": _render_icon_grid,
    "three_column": _render_three_column,
    "comparison": _render_comparison,
    "process_flow": _render_process_flow,
    "timeline": _render_timeline,
    "summary": _render_summary,
    "closing": _render_closing,
}


def export_pptx(
    slide_spec: dict,
    slide_contents: list[dict],
    slide_designs: list[dict],
) -> io.BytesIO:
    """Generate a .pptx file from pipeline output.

    Args:
        slide_spec: The ppt_state spec with presentation metadata and slide list.
        slide_contents: List of {slide_id, type, content} from text_generator.
        slide_designs: List of {slide_id, type, image_b64} from design_generator.

    Returns:
        BytesIO buffer containing the .pptx file.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank

    # Build lookup maps
    ppt_state = slide_spec.get("ppt_state", slide_spec)
    presentation = ppt_state.get("presentation", {})
    meta = presentation.get("meta", {})
    theme = meta.get("theme", {})
    slides = presentation.get("slides", [])

    # Merge defaults
    for key, val in DEFAULT_THEME.items():
        theme.setdefault(key, val)

    bg_color = theme.get("background", "#F5F7FA")
    contents_map = {c["slide_id"]: c.get("content", {}) for c in slide_contents}

    logger.info("[PPTX] Exporting %d slides from slide_spec", len(slides))
    logger.info("[PPTX] slide_contents has %d entries: %s",
                len(slide_contents),
                [c.get("slide_id") for c in slide_contents])
    for s in slides:
        has_content = bool(s.get("content"))
        has_in_map = s["slide_id"] in contents_map
        logger.info("[PPTX]   slide %s type=%s content_in_spec=%s content_in_map=%s",
                    s["slide_id"], s["type"], has_content, has_in_map)

    total_slides = len(slides)

    for slide_idx, slide_data in enumerate(slides):
        slide_id = slide_data["slide_id"]
        slide_type = slide_data["type"]
        content = slide_data.get("content") or contents_map.get(slide_id, {})

        slide = prs.slides.add_slide(blank_layout)

        # Set background color (matching Preview's THEME.background)
        _set_slide_background(slide, bg_color)

        # Render with type-specific renderer
        renderer = RENDERERS.get(slide_type)
        if renderer:
            try:
                renderer(slide, content, theme)
            except Exception as e:
                logger.warning("Failed to render %s (%s): %s", slide_id, slide_type, e)
                _add_textbox(
                    slide, 1.0, 3.0, 11.0, 1.0,
                    content.get("title", slide_id),
                    font_size=24,
                    color=theme.get("text_color", "#1A202C"),
                    bold=True,
                    alignment=PP_ALIGN.CENTER,
                )
        else:
            # Generic fallback: title + any description
            _add_textbox(
                slide, 1.0, 2.5, 11.3, 1.0,
                content.get("title", slide_id),
                font_size=28,
                color=theme.get("text_color", "#1A202C"),
                bold=True,
                alignment=PP_ALIGN.CENTER,
            )
            if content.get("description"):
                _add_textbox(
                    slide, 1.5, 3.8, 10.3, 1.5,
                    content["description"],
                    font_size=16,
                    color=TEXT_SECONDARY,
                    alignment=PP_ALIGN.CENTER,
                )

        # Add slide number to every slide
        _add_slide_number(slide, slide_idx + 1, total_slides)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def export_pptx_from_images(
    slide_spec: dict,
    slide_images: list[str],
) -> io.BytesIO:
    """Generate a .pptx file using captured Preview screenshots.

    Each slide image is placed as a full-bleed background, producing
    output identical to the Preview rendering.

    Args:
        slide_spec: The ppt_state spec (used for metadata/title).
        slide_images: List of base64-encoded PNG images, one per slide.

    Returns:
        BytesIO buffer containing the .pptx file.
    """
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # Blank

    logger.info("[PPTX] Image-based export: %d slide images", len(slide_images))

    for i, image_b64 in enumerate(slide_images):
        if not image_b64:
            logger.warning("[PPTX] Slide %d has empty image, skipping", i)
            continue

        slide = prs.slides.add_slide(blank_layout)

        try:
            image_bytes = base64.b64decode(image_b64)
            image_stream = io.BytesIO(image_bytes)
            slide.shapes.add_picture(
                image_stream,
                Emu(0), Emu(0),
                SLIDE_WIDTH, SLIDE_HEIGHT,
            )
        except Exception as e:
            logger.warning("[PPTX] Failed to add image for slide %d: %s", i, e)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
